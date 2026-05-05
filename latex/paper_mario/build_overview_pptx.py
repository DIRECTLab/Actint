"""Build slides_overview.pptx — 4-slide general overview deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy
import os

BASE = os.path.dirname(os.path.abspath(__file__))
OUTPUTS = os.path.join(BASE, "..", "outputs")
FIGURES = os.path.join(OUTPUTS, "figures")

# ── colours ──────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x00, 0x33, 0x66)
OCEAN   = RGBColor(0x00, 0x66, 0xB3)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
RED     = RGBColor(0xCC, 0x00, 0x00)
BLACK   = RGBColor(0x1A, 0x1A, 0x1A)
LGREY   = RGBColor(0xF2, 0xF4, 0xF7)
DGREY   = RGBColor(0x55, 0x55, 0x55)

# ── slide size: 16:9 ─────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank


def add_slide():
    return prs.slides.add_slide(blank_layout)


def rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def header(slide, title_text):
    """Navy header bar with white title."""
    bar = rect(slide, 0, 0, W, Inches(0.75), fill=NAVY)
    tf = slide.shapes.add_textbox(Inches(0.3), Inches(0.08), W - Inches(0.6), Inches(0.6))
    p = tf.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = WHITE


def textbox(slide, text, x, y, w, h, size=11, bold=False, color=BLACK,
            align=PP_ALIGN.LEFT, wrap=True):
    tf_shape = slide.shapes.add_textbox(x, y, w, h)
    tf = tf_shape.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf_shape


def para(tf, text, size=11, bold=False, color=BLACK, indent=False, space_before=6):
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    p = tf.text_frame.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_before = _Pt(space_before)
    if indent:
        p.level = 1
    run = p.add_run()
    run.text = text
    run.font.size = _Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def image(slide, path, x, y, w, h=None):
    if h:
        return slide.shapes.add_picture(path, x, y, w, h)
    return slide.shapes.add_picture(path, x, y, w)


def caption(slide, text, x, y, w):
    tf = slide.shapes.add_textbox(x, y, w, Inches(0.4))
    tf.text_frame.word_wrap = True
    p = tf.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.italic = True
    run.font.color.rgb = DGREY


def bullet_block(slide, items, x, y, w, h, title=None, size=10.5):
    """Render a list of (text, is_bullet) tuples into a textbox."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        first = False
        run = p.add_run()
        run.text = title
        run.font.size = Pt(size + 0.5)
        run.font.bold = True
        run.font.color.rgb = NAVY
        p.space_before = Pt(0)

    for text, bullet in items:
        if first and not title:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4 if bullet else 8)
        if bullet:
            p.level = 1
        run = p.add_run()
        run.text = ("• " if bullet else "") + text
        run.font.size = Pt(size)
        run.font.color.rgb = BLACK if bullet else NAVY
        run.font.bold = not bullet


def table(slide, headers, rows, x, y, w, col_widths=None, title=None):
    """Simple bordered table."""
    n_cols = len(headers)
    n_rows = len(rows) + 1
    row_h  = Inches(0.33)
    total_h = row_h * n_rows

    if title:
        tb = slide.shapes.add_textbox(x, y - Inches(0.3), w, Inches(0.28))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = NAVY

    tbl = slide.shapes.add_table(n_rows, n_cols, x, y, w, total_h).table

    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw

    def cell_style(cell, text, bold=False, bg=None, align=PP_ALIGN.LEFT, size=10):
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.alignment = align
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = WHITE if bg == NAVY else BLACK
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LGREY if bold else WHITE

    for j, h in enumerate(headers):
        cell_style(tbl.cell(0, j), h, bold=True, bg=NAVY, align=PP_ALIGN.CENTER)

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell_style(tbl.cell(i + 1, j), val,
                       bg=LGREY if i % 2 == 0 else WHITE,
                       align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — The Problem
# ─────────────────────────────────────────────────────────────────────────────
s1 = add_slide()
rect(s1, 0, 0, W, H, fill=WHITE)
header(s1, "The Problem: Vessels Go Dark and Disguise Their Activity")

LX = Inches(0.35)
RX = Inches(6.9)
CY = Inches(0.85)
CW = Inches(6.3)
CH = H - Inches(1.0)

bullet_block(s1, [
    ("90,000 vessels of regulatory concern active at any moment — the picture is incomplete by design.", False),
    ("", False),
    ("75% of industrial fishing vessels intentionally disable their AIS transponders.", False),
    ("", False),
    ("Even when a vessel is visible, analysts must answer two questions from limited data:", False),
    ("What is it doing?  Fishing · Transiting · Loitering · STS transfer", True),
    ("What type is it?  Trawler · Cargo · Naval · Unknown contact", True),
    ("", False),
    ("Available data ranges from a full day of AIS pings to a single satellite image.", False),
], LX, CY, CW, CH, size=11.5)

img_b2 = os.path.join(FIGURES, "B2_dark_period_uncertainty_cones.png")
img_c3 = os.path.join(FIGURES, "C3_sts_rendezvous.png")
image(s1, img_b2, RX, Inches(0.85), CW, Inches(2.9))
image(s1, img_c3, RX, Inches(3.9),  CW, Inches(2.9))
caption(s1, "Top: position uncertainty grows after AIS loss.  Bottom: simulated ship-to-ship transfer.",
        RX, Inches(6.85), CW)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — The Scale (GFW heatmap)
# ─────────────────────────────────────────────────────────────────────────────
s2 = add_slide()
rect(s2, 0, 0, W, H, fill=WHITE)
header(s2, "The Scale: Where Fishing Actually Happens")

bullet_block(s2, [
    ("Grounded in real fleet data from Global Fishing Watch — 2023 effort across four operational regions.", False),
    ("", False),
    ("Brazil EEZ          309,711 records · 1,382,042 fishing hours", True),
    ("Philippine EEZ      61,460 records", True),
    ("Strait of Malacca   54,040 records", True),
    ("Gulf of Guinea      59,954 records", True),
    ("", False),
    ("Fishing effort clusters around productive grounds and shifts seasonally.", False),
    ("", False),
    ("The simulator is calibrated against these distributions — speed profiles, gear-type ratios, and dark-event rates are realistic.", False),
    ("", False),
    ("Validation: speed profile KL-divergence vs. GFW real data < 0.15 nats across all gear types.", False),
], LX, CY, CW, CH, size=11.5)

img_b1 = os.path.join(FIGURES, "B1_gfw_fishing_effort_asia_pacific.png")
image(s2, img_b1, RX, Inches(0.85), CW, Inches(5.9))
caption(s2,
        "GFW 2023 fishing effort — Asia-Pacific. Dense clusters mark the Philippine EEZ and South China Sea. "
        "The Strait of Malacca is a key chokepoint for both fishing and transit traffic.",
        RX, Inches(6.85), CW)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — The System
# ─────────────────────────────────────────────────────────────────────────────
s3 = add_slide()
rect(s3, 0, 0, W, H, fill=WHITE)
header(s3, "The System: Classify Any Vessel from Any Observation")

bullet_block(s3, [
    ("Inputs:", False),
    ("AIS, radar, satellite EO/SAR — or any combination, down to a single ping.", True),
    ("", False),
    ("Outputs per vessel:", False),
    ("Activity class  (fishing / transit / anchored / loitering / STS)", True),
    ("Vessel type  (12 classes)", True),
    ("Dark-vessel risk score  (0–1, SHAP-interpretable)", True),
    ("Calibrated confidence  (conservative for sparse data)", True),
    ("", False),
    ("Key design insight:", False),
    ("XGBoost learns what each feature's absence means — no imputation, no minimum track length required.", True),
    ("", False),
    ("Trained and evaluated on 5,121 real vessels across three independent NOAA datasets.", False),
    ("Training time: 82 seconds on GPU.   Inference: 19 µs per vessel.", False),
], LX, CY, CW, CH, size=11)

img_dash = os.path.join(FIGURES, "B5_risk_score_dashboard.png")
img_feat = os.path.join(OUTPUTS, "brazil_eez", "feature_importance.png")
image(s3, img_dash, RX, Inches(0.85), CW, Inches(3.2))
image(s3, img_feat, RX, Inches(4.15), CW, Inches(2.6))
caption(s3, "Top: risk dashboard — all flagged vessels are fishing-registered.  Bottom: top features by importance.",
        RX, Inches(6.85), CW)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Results
# ─────────────────────────────────────────────────────────────────────────────
s4 = add_slide()
rect(s4, 0, 0, W, H, fill=WHITE)
header(s4, "Results: Robust Performance Across Regions and Sensors")

col_w = [Inches(3.5), Inches(1.8)]
table(s4,
      ["Metric", "Result"],
      [
          ["Activity classification F₁",        "0.900 – 1.000"],
          ["Vessel type F₁  (12 classes)",       "0.969"],
          ["Performance: N=1 vs N=50",           "statistically flat"],
          ["Brazil EEZ  (synthetic)",            "F₁ = 0.993"],
          ["Strait of Malacca  (synthetic)",     "F₁ = 0.999"],
          ["Gulf of Guinea  (synthetic)",        "F₁ = 0.992"],
          ["Philippine EEZ  (synthetic)",        "F₁ = 0.992"],
          ["Dark vessels flagged",               "49  (all fishing-reg.)"],
          ["GPU training  (216k examples)",      "82 seconds"],
      ],
      LX, Inches(0.9), Inches(5.5), col_widths=col_w,
)

tb = slide = s4
note = s4.shapes.add_textbox(LX, Inches(4.8), Inches(5.5), Inches(0.55))
note.text_frame.word_wrap = True
p = note.text_frame.paragraphs[0]
run = p.add_run()
run.text = ("The most significant open challenge is EO/SAR-only contacts where speed is unavailable. "
            "Geography-prior features and SAR foundation models are the next step.")
run.font.size = Pt(10)
run.font.italic = True
run.font.color.rgb = DGREY

img_f1   = os.path.join(FIGURES, "A2_partial_track_f1_vs_length.png")
img_conf = os.path.join(OUTPUTS, "brazil_eez", "confusion_matrix.png")
image(s4, img_f1,   RX, Inches(0.85), CW, Inches(2.9))
image(s4, img_conf, RX, Inches(3.9),  CW, Inches(2.9))
caption(s4, "Top: F₁ is flat from N=1 to N=50 — no minimum custody duration needed.  Bottom: confusion matrix, Brazil EEZ.",
        RX, Inches(6.85), CW)


# ─────────────────────────────────────────────────────────────────────────────
out = os.path.join(BASE, "slides_overview.pptx")
prs.save(out)
print(f"Saved {out}")
